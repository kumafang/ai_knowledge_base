#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Fri Jan  9 19:20:00 2026

@author: kumafang
"""

import pdfplumber
import pandas as pd
from docx import Document
from pptx import Presentation

def load_file(file):
    name = file.name.lower()

    if name.endswith(".pdf"):
        with pdfplumber.open(file) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)

    if name.endswith(".csv"):
        df = pd.read_csv(file)
        return df.to_string()

    if name.endswith(".xlsx"):
        df = pd.read_excel(file)
        return df.to_string()

    if name.endswith(".docx"):
        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    if name.endswith(".pptx"):
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    return ""
